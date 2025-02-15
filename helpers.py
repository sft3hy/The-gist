import streamlit as st
import base64
import os
from PyPDF2 import PdfReader
from pptx import Presentation


def get_base64_of_bin_file(png_file: str) -> str:
    with open(png_file, "rb") as f:
        return base64.b64encode(f.read()).decode()

@st.cache_resource
def build_markup_for_logo(png_file: str) -> str:
    # Convert the PNG file to a base64 encoded binary string
    binary_string = get_base64_of_bin_file(png_file)
    # Create a string of HTML code that sets the background image of a sidebar header
    return f"""
            <style>
                [data-testid="stSidebarHeader"] {{
                    background-image: url("data:image/png;base64,{binary_string}");
                    background-repeat: no-repeat;
                    background-size: 200px;
                    padding-bottom: 5rem;
                    background-position: top center;
                }}
            </style>
            """

def extract_text(file_path):
    """
    Extracts text from a PDF, PPTX, or TXT file.

    Args:
        file_path (str): Path to the file.

    Returns:
        str: Extracted text from the file.
    """
    # Check if the file exists
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"The file {file_path} does not exist.")
    
    # Get the file extension
    file_extension = os.path.splitext(file_path)[1].lower()

    potential_already_extracted = f"text_extractions/{file_path.split('/')[-1]}.txt"
    if os.path.exists(potential_already_extracted):
        return open(potential_already_extracted, 'r').read()
    else:
    # Extract text based on file type
        if file_extension == '.pdf':
            return extract_text_from_pdf(file_path)
        elif file_extension == '.pptx':
            return extract_text_from_pptx(file_path)
        elif file_extension == '.txt':
            return extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

# Define a function to extract text from a PDF file
def extract_text_from_pdf(file_path):
    # Initialize an empty string to store the extracted text
    text = ""
    # Open the PDF file in binary read mode
    with open(file_path, 'rb') as pdf_file:
        # Create a PdfReader object to read the PDF file
        reader = PdfReader(pdf_file)
        # Iterate over each page in the PDF file
        for page in reader.pages:
            # Extract the text from the current page and append it to the text string
            text += page.extract_text() + "\n"
    # Write the extracted text to a file in the extractions folder
    write_text_to_extractions_folder(text, f"text_extractions/{file_path.split('/')[-1]}.txt")
    # Return the extracted text
    return text

# Define a function to extract text from a PPTX file
def extract_text_from_pptx(file_path):
    # Initialize an empty string to store the extracted text
    text = ""
    # Open the PPTX file
    presentation = Presentation(file_path)
    # Iterate over each slide in the presentation
    for slide in presentation.slides:
        # Iterate over each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has a text frame
            if shape.has_text_frame:
                # Extract the text from the shape and append it to the text string
                text += shape.text + "\n"
    # Write the extracted text to a file in the extractions folder
    write_text_to_extractions_folder(text, f"text_extractions/{file_path.split('/')[-1]}.txt")
    # Return the extracted text
    return text

# Define a function to extract text from a TXT file
def extract_text_from_txt(file_path):
    # Open the TXT file in read mode with UTF-8 encoding
    with open(file_path, 'r', encoding='utf-8') as txt_file:
        # Read the text from the file
        text = txt_file.read()
        # Write the extracted text to a file in the extractions folder
        write_text_to_extractions_folder(text, f"text_extractions/{file_path.split('/')[-1]}.txt")
        # Return the extracted text
        return txt_file.read()

# Define a function to write text to a file in the extractions folder
def write_text_to_extractions_folder(text, file_path):
    # Extract the file name from the file path
    file_name = file_path.split('/')[-1]
    # Construct the new file path
    new_path = f"text_extractions/{file_name}"
    # Open the new file in write mode
    with open(new_path, 'w') as f:
        # Write the text to the file
        f.write(text)
